from pptx import Presentation # 라이브러리 
from pptx.util import Inches # 사진, 표등을 그리기 위해

prs = Presentation()

with open('lyrics.txt', 'r', encoding='utf-8') as f:
    while True:
        rst = []
        line = f.readline()
        if line == '\n':
            continue
        if not line: break

        rst.append(line)

        line = f.readline()

        if not line: break

        rst.append(line) 

        title_slide_layout = prs.slide_layouts[5] # 5 : 제목슬라이드에 해당
        slide = prs.slides.add_slide(title_slide_layout) # 슬라이드 추가

        # 제목 - 제목에 값넣기
        title = slide.placeholders[0] # 제목
        title.text = rst[0] + rst[1] # 제목에 값 넣기
# 저장
prs.save('lyrics.pptx')