from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN   # 정렬 설정하기
from pptx.util import Pt      # Pt 폰트사이즈

test_name = './test.pptx'
prs = Presentation(test_name)

slide_num = 0
slide = prs.slides[slide_num]

shape_name = 'name1'

# 슬라이드 내 shape 사전 만들기
shapes_list = slide.shapes
shape_index = {}
for i, shape in enumerate(shapes_list):
    shape_index[ shape.name ] = i
print(shape_index)   # {'Box_down': 0, 'Box_up': 1, 'name2': 2, 'name1': 3}

def text_on_shape(shape, input_text, font_size = 95, bold = True):

    # shape 내 텍스트 프레임 선택하기 & 기존 값 삭제하기
    text_frame = shape.text_frame
    text_frame.clear()

    # 문단 선택하기
    p = text_frame.paragraphs[0]

    # 정렬 설정 : 중간정렬
    p.alighnment = PP_ALIGN.CENTER   

    # 텍스트 입력 / 폰트 지정
    run = p.add_run()
    run.text = input_text
    font = run.font
    font.size = Pt(font_size)
    font.bold = bold
    font.name = None  # 지정하지 않으면 기본 글씨체로  #  'Arial'

# 하나의 shape 선택하기
shape_name = 'name1'
shape_select = shapes_list[ shape_index[ shape_name ]]    

text_on_shape(shape_select, '홍길동')

prs.save('test2.pptx')